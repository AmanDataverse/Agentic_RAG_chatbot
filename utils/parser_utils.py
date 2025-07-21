import os
import fitz  # PyMuPDF
import docx
import pandas as pd
from pptx import Presentation

def parse_document(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".csv":
        return parse_csv(file_path)
    elif ext in [".txt", ".md"]:
        return parse_text(file_path)
    else:
        return ["Unsupported file type."]

def parse_pdf(path):
    doc = fitz.open(path)
    text = ""
    for page in doc:
        text += page.get_text()
    return chunk_text(text)

def parse_docx(path):
    doc = docx.Document(path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return chunk_text(text)

def parse_pptx(path):
    prs = Presentation(path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return chunk_text(text)

def parse_csv(path):
    df = pd.read_csv(path)
    return chunk_text(df.to_string())

def parse_text(path):
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
    return chunk_text(text)

def chunk_text(text, max_len=500):
    sentences = text.split(". ")
    chunks, chunk = [], ""

    for sentence in sentences:
        if len(chunk) + len(sentence) < max_len:
            chunk += sentence + ". "
        else:
            chunks.append(chunk.strip())
            chunk = sentence + ". "

    if chunk:
        chunks.append(chunk.strip())

    return chunks
